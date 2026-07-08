import os

os.environ.setdefault("DASHSCOPE_API_KEY", "")

import asyncio
from typing import Annotated, TypedDict
from operator import add
from pathlib import Path
from langchain_classic.agents import create_react_agent
from langchain_community.chat_models import ChatTongyi
from langchain_community.embeddings import DashScopeEmbeddings
from langchain_core.messages import AnyMessage, HumanMessage
from langgraph.checkpoint.memory import InMemorySaver
from langgraph.graph import StateGraph
from langgraph.constants import END, START
from core.DocProcess.document_process import doc_processor
from core.DrawingRecoAssistant.drawing_assistant import get_drawing_assistant
from settings.Define import Params, PathConfig
from settings.logger_manager import get_logger

try:
    from langchain_mcp_adapters.client import MultiServerMCPClient
except ImportError:
    MultiServerMCPClient = None

logger = get_logger(__name__)

llm = ChatTongyi(
    model=Params.DEFAULT_CHAT_MODEL,
    api_key=os.getenv("DASHSCOPE_API_KEY"),
    api_base=Params.API_BASE
)

embedding_model = DashScopeEmbeddings(model=Params.DEFAULT_EMBEDDING_MODEL)

class State(TypedDict):
    messages: Annotated[list[AnyMessage], add]
    type: str
    file_path: str  # 添加文件路径字段
    previous_node: str  # 记录上一个执行的节点类型，用于反思后路由
    skill: str  # 前端选择的技能类型：travel/joke/couplet/document/code/drawing/other


def create_agent_graph():
    nodes = Params.NODE_LIST

    def other_node(state: State):
        logger.info("other_node 开始处理")
        return {"messages": [HumanMessage(content="抱歉，您咨询的问题不在我的能力范围内，无法回答这类问题。")],
                "type": "other", "previous_node": "other_node"}

    def supervisor_node(state: State):
        logger.info("supervisor_node")
        prompts = """你是一个专业的客服助手，需要根据用户的问题进行任务分类，并将任务分给对应的Agent来执行。

分类规则：
1. 如果用户的问题是和旅游路线规划相关的，那就返回travel。
2. 如果用户的问题是和讲笑话相关的，那就返回joke。
3. 如果用户的问题是和对联相关的，那就返回couplet。
4. 如果用户的问题是和文档相关的（包含修改已有文档、新建文档并写入内容、排版、修正格式、修改对齐方式，查看文档中的错别字等等），那就返回document。
5. 如果用户的问题是和编程相关的，那就返回code。
6. 如果用户的问题是和工程图纸识别、尺寸提取、公差校核、DXF解析、国标查询、单位换算等相关的，那就返回drawing。
7. 如果是其它问题，那就返回other。

重要：如果用户的问题是追问、确认、或者对之前回答的反馈（例如"确认下是否写得对？"、"再优化一下"、"不对，应该是..."），请根据对话历史判断之前讨论的主题，并返回对应的分类。
除了这几个选项外，不要返回任何其它的内容。
"""

        message = state["messages"][-1]
        if hasattr(message, 'content'):
            message_content = message.content
        else:
            message_content = str(message)
        logger.info(f"用户消息: {message_content[:30]}...")

        # Build conversation history for context
        conversation_history = []
        for msg in state["messages"][:-1]:
            if hasattr(msg, 'content'):
                content = msg.content
                msg_type = type(msg).__name__
            else:
                content = str(msg)
                msg_type = "dict"

            if isinstance(msg, dict) and msg.get("role") == "assistant":
                conversation_history.append({"role": "assistant", "content": msg["content"]})
            elif "HumanMessage" in str(type(msg)) or (isinstance(msg, dict) and msg.get("role") == "user"):
                conversation_history.append({"role": "user", "content": content})

        prompt_list = [
            {"role": "system", "content": prompts},
        ] + conversation_history + [
            {"role": "user", "content": message_content},
        ]
        if "type" in state:
            logger.info("已有类型，返回 END")
            return {"type": END}
        elif state.get("skill", "").strip() in nodes:
            skill_type = state["skill"].strip()
            logger.info(f"前端指定技能: {skill_type}，直接路由")
            return {"type": skill_type}
        else:
            try:
                response = llm.invoke(prompt_list)
                typeRes = response.content.strip()
                logger.debug(f"分类结果原始值: '{typeRes}' (长度: {len(typeRes)})")
                logger.debug(f"分类结果.strip(): '{typeRes.strip()}'")
                if typeRes.strip() in nodes:
                    logger.info(f"分类成功: {typeRes.strip()}")
                    return {"type": typeRes.strip()}
                else:
                    logger.warning(f"分类失败: '{typeRes}' 不在列表 {nodes} 中")
                    logger.info("默认返回 other")
                    return {"type": "other"}
            except Exception as e:
                logger.error(f"supervisor_node LLM 调用失败: {str(e)[:50]}")
                return {"type": "other"}

    def travel_node(state: State):
        system_prompt = """你是一个专业的旅行规划助手，根据用户的问题，生成一个旅游路线规划。请用中文回答，并返回一个不超过100字的规划结果"""

        # Get the last message and build conversation history
        message = state["messages"][-1]
        if hasattr(message, 'content'):
            message_content = message.content
        else:
            message_content = str(message)

        conversation_history = []
        for msg in state["messages"][:-1]:
            if hasattr(msg, 'content'):
                content = msg.content
            else:
                content = str(msg)

            if isinstance(msg, dict) and msg.get("role") == "assistant":
                conversation_history.append({"role": "assistant", "content": content})
            elif "HumanMessage" in str(type(msg)) or (isinstance(msg, dict) and msg.get("role") == "user"):
                conversation_history.append({"role": "user", "content": content})

        prompt_list = [
            {"role": "system", "content": system_prompt},
        ] + conversation_history + [
            {"role": "user", "content": message_content},
        ]
        try:
            if MultiServerMCPClient:
                client = MultiServerMCPClient(
                    {"amap-maps-sse": {"url": "https://mcp.amap.com/sse?key=77067076328eb3b90559c0cb1222ff81",
                                       "transport": "streamable_http"}}
                )
                tools = asyncio.run(client.get_tools())
                agent = create_react_agent(model=llm, tools=tools)
                response = agent.invoke({"messages": prompt_list})
            else:
                raise ImportError("MultiServerMCPClient not available")
        except Exception as e:
            logger.warning(f"MCP 服务不可用，使用本地规划: {str(e)[:50]}")
            response = llm.invoke(prompt_list)
        content = response.content if hasattr(response, 'content') else str(response)
        if not content or content.strip() == "":
            content = "抱歉，暂时无法生成旅游路线规划。"
        return {"messages": [HumanMessage(content=content)], "type": "travel", "previous_node": "travel_node"}

    def joke_node(state: State):
        logger.info("joke_node")
        system_prompt = """你是一个专业的笑话大师，根据用户的问题，写一个不超过100个字的笑话。"""

        message = state["messages"][-1]
        if hasattr(message, 'content'):
            message_content = message.content
        else:
            message_content = str(message)

        conversation_history = []
        for msg in state["messages"][:-1]:
            if hasattr(msg, 'content'):
                content = msg.content
            else:
                content = str(msg)

            if isinstance(msg, dict) and msg.get("role") == "assistant":
                conversation_history.append({"role": "assistant", "content": content})
            elif "HumanMessage" in str(type(msg)) or (isinstance(msg, dict) and msg.get("role") == "user"):
                conversation_history.append({"role": "user", "content": content})

        prompt_list = [
            {"role": "system", "content": system_prompt},
        ] + conversation_history + [
            {"role": "user", "content": message_content},
        ]
        try:
            response = llm.invoke(prompt_list)
            content = response.content if hasattr(response, 'content') else str(response)
        except Exception as e:
            logger.error(f"joke_node LLM 调用失败: {str(e)[:50]}")
            content = "抱歉，当前服务暂时不可用，请稍后重试。"
        return {"messages": [HumanMessage(content=content)], "type": "joke", "previous_node": "joke_node"}

    def code_node(state: State):
        logger.info("code_node")

        # 获取文件名
        file_path = state.get("file_path")
        filename = "代码"
        if file_path:
            from pathlib import Path
            filename = Path(file_path).name

        system_prompt = f"""你是一个专业的代码编程专家，熟悉各种语言的编程，比如C、C++、HTML、Go、Java、PHP、Python等。

请根据用户的具体需求进行回答：

**场景A：用户要求生成新代码**
- 使用 Markdown 格式，包括标题（# ## ###）、加粗（**文本**）、列表（- 或 1.）等
- 代码块必须独占成段，使用三个反引号包裹
- 代码要完整，包含必要的导入、函数定义、测试用例和运行示例
- 代码要有注释，关键逻辑处添加简洁的中文注释
- 分版本展示：如果适用，用二级标题（##）分隔不同版本，如"## 1. 基础版"和"## 2. 优化版"
- 复杂度分析：在代码后用列表形式简要说明时间复杂度和空间复杂度
- 运行示例：给出代码的运行输出示例

回答结构示例（生成新代码）：
## 标题

简要说明文字。

### 1. 基础版

核心逻辑说明。

```python
# 完整代码
```

### 2. 优化版

优化说明。

```python
# 完整代码
```

### 复杂度分析

- **时间复杂度**：说明
- **空间复杂度**：说明

### 运行示例

```
输出结果
```

**场景B：用户要求审查/修改已有代码**
- 仔细检查代码中的语法问题、逻辑问题、BUG问题等
- 如果没有发现任何问题，必须回复："{filename}代码文件不需要修改，没有发现任何明显性错误"
- 如果发现问题，必须按照以下格式回复：

发现{{N}}处问题：
第1处：{{问题描述}}
第2处：{{问题描述}}
...

修改后的完整代码：
```{{language}}
{{修改后的完整代码}}
```

修改后的文件下载路径：http://127.0.0.1:5001/download/{{output_filename}}

注意：代码块的语言标识必须正确（如 python、javascript、java、cpp 等），以便前端正确高亮显示。"""

        message = state["messages"][-1]
        if hasattr(message, 'content'):
            message_content = message.content
        else:
            message_content = str(message)

        conversation_history = []
        for msg in state["messages"][:-1]:
            if hasattr(msg, 'content'):
                content = msg.content
            else:
                content = str(msg)

            if isinstance(msg, dict) and msg.get("role") == "assistant":
                conversation_history.append({"role": "assistant", "content": content})
            elif "HumanMessage" in str(type(msg)) or (isinstance(msg, dict) and msg.get("role") == "user"):
                conversation_history.append({"role": "user", "content": content})

        prompt_list = [
            {"role": "system", "content": system_prompt},
        ] + conversation_history + [
            {"role": "user", "content": message_content},
        ]
        try:
            response = llm.invoke(prompt_list)
            content = response.content if hasattr(response, 'content') else str(response)
        except Exception as e:
            logger.error(f"code_node LLM 调用失败: {str(e)[:50]}")
            content = "抱歉，当前服务暂时不可用，请稍后重试。"

        # 检查是否需要保存修改后的代码文件
        if file_path and "修改后的完整代码" in content:
            try:
                from pathlib import Path
                import uuid

                # 提取代码块内容
                import re
                code_match = re.search(r'```(?:\w+)?\n([\s\S]*?)```', content)
                if code_match:
                    modified_code = code_match.group(1).strip()

                    # 获取原始文件扩展名和文件名
                    original_ext = Path(file_path).suffix
                    original_name = Path(file_path).stem

                    # 生成输出文件名
                    output_filename = f"modified_{original_name}_{str(uuid.uuid4())[:8]}{original_ext}"
                    output_path = PathConfig.OUTPUTS_DIR / output_filename

                    # 保存修改后的代码
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(modified_code)

                    # 生成下载链接
                    download_url = f"http://127.0.0.1:5001/download/{output_filename}"

                    # 替换内容中的下载路径（支持占位符或LLM生成的原始文件名链接）
                    content = content.replace("http://127.0.0.1:5001/download/{output_filename}", download_url)
                    # 也替换LLM可能生成的原始文件名链接
                    original_download_pattern = rf"http://127\.0\.0\.1:5001/download/{re.escape(Path(file_path).name)}"
                    content = re.sub(original_download_pattern, download_url, content)

                    logger.info(f"代码文件保存成功: {output_path}")
            except Exception as e:
                logger.error(f"保存代码文件失败: {str(e)}")

        return {"messages": [HumanMessage(content=content)], "type": "code", "previous_node": "code_node"}

    def couplet_node(state: State):
        logger.info("couplet_node 开始处理")

        from core.Couplet.couplet import Couplet

        couplet = Couplet()
        if not couplet.load_vector():
            return {"messages": [HumanMessage(content="抱歉，向量数据库加载失败。")], "type": "couplet"}


        query = state["messages"][-1]
        if hasattr(query, 'content'):
            query_text = query.content
        else:
            query_text = str(query)
        logger.info(f"查询文本: {query_text}")

        prompt = couplet.similarity_search(query_text)
        if not prompt:
            return {"messages": [HumanMessage(content="抱歉，Prompt 构建失败。")], "type": "couplet"}

        try:
            # RAG生成
            response = llm.invoke(prompt)
            logger.info("LLM 调用成功")
        except Exception as e:
            logger.error(f"LLM 调用失败: {str(e)}")
            return {"messages": [HumanMessage(content="抱歉，LLM 调用失败。")], "type": "couplet"}

        content = response.content if hasattr(response, 'content') else str(response)
        if not content or content.strip() == "":
            content = "抱歉，暂时无法生成对联。"

        logger.info(f"对联结果: {content[:30]}...")
        return {"messages": [HumanMessage(content=content)], "type": "couplet", "previous_node": "couplet_node"}

    def document_node(state: State):
        logger.info("document_node 开始处理")

        message = state["messages"][-1]
        if hasattr(message, 'content'):
            message_content = message.content
        else:
            message_content = str(message)

        # Build conversation history for context
        conversation_history = []
        for msg in state["messages"][:-1]:
            if hasattr(msg, 'content'):
                content = msg.content
            else:
                content = str(msg)

            if isinstance(msg, dict) and msg.get("role") == "assistant":
                conversation_history.append({"role": "assistant", "content": content})
            elif "HumanMessage" in str(type(msg)) or (isinstance(msg, dict) and msg.get("role") == "user"):
                conversation_history.append({"role": "user", "content": content})

        logger.info(f"用户输入内容长度: {len(message_content)}")

        file_path = state.get("file_path")

        # 如果没有文件路径，说明是新建文档
        if not file_path:
            return _handle_document_creation(state, message_content, conversation_history)

        logger.info(f"文件路径: {file_path}")

        from pathlib import Path
        file_ext = Path(file_path).suffix.lower()
        logger.info(f"文件扩展名: {file_ext}")

        # 解析用户的修改要求
        # 消息格式：[文档内容]\n...\n\n[用户要求]\n修改要求
        if '[用户要求]' in message_content:
            parts = message_content.split('[用户要求]')
            if len(parts) > 1:
                modification_request = parts[-1].strip()
            else:
                modification_request = "修改文档"
        else:
            lines = message_content.split('\n')
            modification_request = lines[0].strip() if lines else "修改文档"

        logger.info(f"修改要求: {modification_request}")

        # 根据修改要求构建不同的 prompt
        if '错别字' in modification_request or '错字' in modification_request or '错误字' in modification_request:
            system_prompt = """你是一个专业的文档校对助手。你的任务是检查并修正文档中的错别字和词语搭配错误。

    重要要求：
    1. 只返回修改后的文档内容，不要任何解释、说明或额外文字
    2. 保持原文档的换行结构和格式完全一致
    3. 检查词语搭配是否合理（如"天弃"应改为"天气"，"银绗"应改为"银行"）
    4. 修正所有不合理的词语组合和错别字
    5. 如果没有发现错别字，直接回复"没有发现错别字"

    注意：你的回复将直接保存为文档文件，所以只能包含修改后的文档内容本身！"""
        elif '格式' in modification_request or '排版' in modification_request or '对齐' in modification_request:
            is_alignment_request = '对齐' in modification_request

            if is_alignment_request and file_ext == '.txt':
                system_prompt = """你是一个专业的文档格式调整助手。你的任务是根据用户要求调整文档的对齐方式。

    重要要求：
    1. 只调整对齐方式（如左对齐、居中对齐、右对齐），不要修改文字内容
    2. 对于 .txt 纯文本文件：
       - 左对齐：去除每行开头的所有空格和制表符，让内容紧贴左侧
       - 居中对齐：在每行前后添加适当空格使内容居中
       - 右对齐：在每行开头添加适当空格使内容靠右
    3. 保留原文档的所有文字内容
    4. 直接返回修改后的文本，保持原文本的换行结构

    注意：只调整对齐方式，不要改变任何文字内容！"""
            else:
                system_prompt = """你是一个专业的文档格式调整助手。你的任务是根据用户要求调整文档格式。

    要求：
    1. 只调整格式（排版、对齐方式、字体样式等）
    2. 不要修改文档的文字内容
    3. 保留原文档的所有内容
    4. 直接返回修改后的文本，保持原文本的换行结构

    注意：只调整格式，不要改变文字内容。"""
        elif '内容' in modification_request or '修改' in modification_request or '增删' in modification_request:
            system_prompt = """你是一个专业的文档修改助手。你的任务是根据用户要求修改文档内容。

    要求：
    1. 只修改用户指定的内容部分
    2. 保留文档的其他内容不变
    3. 保留原文档的结构和格式
    4. 直接返回修改后的文本，保持原文本的换行结构

    注意：只修改用户要求的部分，不要影响其他内容。"""
        elif '结构' in modification_request or '优化' in modification_request:
            system_prompt = """你是一个专业的文档优化助手。你的任务是优化文档结构和表达。

    要求：
    1. 只优化文档结构和表达方式
    2. 保留文档的核心内容
    3. 不要删除重要信息
    4. 直接返回优化后的文本，保持原文本的换行结构

    注意：只优化结构和表达，不要改变核心内容。"""
        else:
            system_prompt = """你是一个专业的文档修改助手。你需要根据用户的要求来修改文档内容。

    要求：
    1. 只修改用户指定的部分
    2. 保留文档的其他内容不变
    3. 保留原文档的结构和格式
    4. 直接返回修改后的文本，保持原文本的换行结构

    注意：严格按照用户的要求进行修改，不要影响其他部分。"""

        try:
            import uuid
            import shutil

            text_only = ""
            text_paragraphs_info = []
            original_content = ""

            if file_ext == '.docx':
                from docx import Document
                doc = Document(file_path)
                for i, para in enumerate(doc.paragraphs):
                    if para.text.strip():
                        text_paragraphs_info.append((i, para.text))
                text_only = '\n'.join([text for _, text in text_paragraphs_info])
                logger.info(f"[docx] 原文档总段落数: {len(doc.paragraphs)}")
                logger.info(f"[docx] 有文本的段落数: {len(text_paragraphs_info)}")

            elif file_ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    text_only = f.read()
                original_content = text_only
                logger.info(f"[txt] 文本长度: {len(text_only)}")

            elif file_ext in ['.xlsx', '.xls']:
                try:
                    import openpyxl
                    if file_ext == '.xlsx':
                        try:
                            wb = openpyxl.load_workbook(file_path, data_only=True)
                            content_parts = []
                            for sheet in wb.sheetnames:
                                ws = wb[sheet]
                                content_parts.append(f"=== {sheet} ===")
                                for row in ws.iter_rows(values_only=True):
                                    row_text = '\t'.join(str(cell) for cell in row if cell is not None)
                                    if row_text.strip():
                                        content_parts.append(row_text)
                            text_only = '\n'.join(content_parts)
                            logger.info(f"[excel] 使用 openpyxl 读取成功，文本长度: {len(text_only)}")
                        except Exception as e:
                            logger.error(f"[excel] openpyxl 读取失败: {str(e)}")
                            try:
                                import xlrd
                                wb = xlrd.open_workbook(file_path)
                                content_parts = []
                                for sheet_idx in range(wb.nsheets):
                                    sheet = wb.sheet_by_index(sheet_idx)
                                    content_parts.append(f"=== {sheet.name} ===")
                                    for row_idx in range(sheet.nrows):
                                        row_values = sheet.row_values(row_idx)
                                        row_text = '\t'.join(str(val) for val in row_values if val != '')
                                        if row_text.strip():
                                            content_parts.append(row_text)
                                text_only = '\n'.join(content_parts)
                                logger.info(f"[excel] 使用 xlrd 读取成功，文本长度: {len(text_only)}")
                            except Exception as e2:
                                logger.error(f"[excel] xlrd 也读取失败: {str(e2)}")
                                raise ValueError(
                                    f"无法读取 Excel 文件，请确认文件格式正确。\n错误信息1: {str(e)}\n错误信息2: {str(e2)}")

                    elif file_ext == '.xls':
                        import xlrd
                        wb = xlrd.open_workbook(file_path)
                        content_parts = []
                        for sheet_idx in range(wb.nsheets):
                            sheet = wb.sheet_by_index(sheet_idx)
                            content_parts.append(f"=== {sheet.name} ===")
                            for row_idx in range(sheet.nrows):
                                row_values = sheet.row_values(row_idx)
                                row_text = '\t'.join(str(val) for val in row_values if val != '')
                                if row_text.strip():
                                    content_parts.append(row_text)
                        text_only = '\n'.join(content_parts)
                        logger.info(f"[excel] 使用 xlrd 读取 .xls 成功，文本长度: {len(text_only)}")

                except ImportError as e:
                    logger.error(f"[excel] 缺少必要的库: {str(e)}")
                    raise ValueError("缺少必要的库，请安装：pip install openpyxl xlrd")
                except Exception as e:
                    logger.error(f"[excel] 读取 Excel 文件失败: {str(e)}")
                    raise ValueError(f"无法读取 Excel 文件，请确认文件格式正确。\n错误信息: {str(e)}")

            elif file_ext in ['.pptx', '.ppt']:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    content_parts = []
                    for i, slide in enumerate(prs.slides, 1):
                        content_parts.append(f"=== 幻灯片 {i} ===")
                        for shape in slide.shapes:
                            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                                if shape.text_frame and hasattr(shape.text_frame, 'text'):
                                    text_content = shape.text_frame.text.strip()
                                    if text_content:
                                        content_parts.append(text_content)
                    text_only = '\n'.join(content_parts)
                    logger.info(f"[pptx] 读取成功，文本长度: {len(text_only)}")
                except Exception as e:
                    logger.error(f"[pptx] 读取 PowerPoint 文件失败: {str(e)}")
                    raise ValueError(f"无法读取 PowerPoint 文件。\n错误信息: {str(e)}")

            elif file_ext == '.doc':
                logger.info("[doc] 检测到旧版 .doc 格式，尝试读取...")
                try:
                    import win32com.client
                    word = win32com.client.Dispatch("Word.Application")
                    word.Visible = False
                    doc_obj = word.Documents.Open(file_path)
                    text_only = doc_obj.Content.Text
                    doc_obj.Close()
                    word.Quit()
                except ImportError:
                    logger.warning("[doc] 无法直接读取 .doc 文件，请安装 pywin32: pip install pywin32")
                    raise ValueError("暂不支持直接读取 .doc 文件，请先转换为 .docx 格式")
                logger.info(f"[doc] 文本长度: {len(text_only)}")
            else:
                raise ValueError(f"不支持的文件格式: {file_ext}")

            logger.info(f"提取的文本长度: {len(text_only)}")

            prompt_list = [
                {"role": "system", "content": system_prompt},
            ] + conversation_history + [
                {"role": "user",
                 "content": f"修改要求：{modification_request}\n\n文档内容：\n{text_only}"},
            ]

            response = llm.invoke(prompt_list)
            content = response.content if hasattr(response, 'content') else str(response)

            if not content or content.strip() == "":
                content = "抱歉，暂时无法处理文档修改请求。"
            else:
                logger.info(f"LLM 返回内容长度: {len(content)}")

                from pathlib import Path
                original_filename = Path(file_path).stem
                output_filename = f"{original_filename}_{uuid.uuid4().hex[:8]}{file_ext}"
                output_path = doc_processor.OUTPUT_DIR / output_filename

                if file_ext == '.docx':
                    from docx import Document
                    doc = Document(file_path)
                    modified_lines = content.split('\n')
                    logger.info(f"修改后行数: {len(modified_lines)}")

                    for idx, (para_index, original_text) in enumerate(text_paragraphs_info):
                        if idx < len(modified_lines):
                            para = doc.paragraphs[para_index]
                            new_text = modified_lines[idx]
                            for run in para.runs:
                                run.text = ""
                            if new_text:
                                para.add_run(new_text)

                    doc.save(str(output_path))

                elif file_ext == '.txt':
                    is_typo_request = '错别字' in modification_request or '错字' in modification_request or '错误字' in modification_request

                    if is_typo_request:
                        no_change_markers = ['没有发现', '未发现', '没有错别字', '无需修改', '原文档内容',
                                             '文档内容正确']
                        is_no_change = any(marker in content for marker in no_change_markers)

                        if is_no_change:
                            logger.info("[txt] LLM 返回说明文字，不保存文档")
                            content = "没有发现错别字"
                            output_path = None
                        else:
                            logger.debug(f"[txt] LLM 返回内容预览: {content[:300]}...")

                            markers = ['修改后的内容', '修改后内容', '以下是修改后的', '修改后的文档', '文档内容如下',
                                       '修改后的文本']
                            extracted_content = None
                            for marker in markers:
                                if marker in content:
                                    idx = content.index(marker)
                                    newline_idx = content.find('\n', idx)
                                    if newline_idx != -1:
                                        extracted_content = content[newline_idx + 1:].strip()
                                        logger.info(
                                            f"[txt] 通过标记 '{marker}' 提取到内容，长度: {len(extracted_content)}")
                                        break

                            if extracted_content is None:
                                if content.startswith('好的') or content.startswith('已经') or content.startswith('我'):
                                    paragraphs = content.split('\n\n')
                                    if len(paragraphs) > 1:
                                        extracted_content = '\n\n'.join(paragraphs[1:]).strip()
                                        logger.info(f"[txt] 通过段落分割提取到内容，长度: {len(extracted_content)}")

                            final_content = extracted_content if extracted_content else content
                            logger.info(f"[txt] 最终保存内容长度: {len(final_content)}")

                            with open(output_path, 'w', encoding='utf-8') as f:
                                f.write(final_content)
                    else:
                        logger.info("[txt] 非错别字请求，直接保存内容")

                        if '对齐' in modification_request or '格式' in modification_request:
                            logger.info("[txt] 检测到对齐/格式请求，执行后处理")
                            lines = content.split('\n')
                            processed_lines = []

                            for line in lines:
                                if '左对齐' in modification_request or '左' in modification_request:
                                    processed_lines.append(line.lstrip())
                                elif '居中' in modification_request:
                                    stripped_line = line.strip()
                                    if stripped_line:
                                        total_spaces = max(0, (80 - len(stripped_line)) // 2)
                                        processed_lines.append(' ' * total_spaces + stripped_line)
                                    else:
                                        processed_lines.append('')
                                elif '右对齐' in modification_request:
                                    stripped_line = line.strip()
                                    if stripped_line:
                                        total_spaces = max(0, 80 - len(stripped_line))
                                        processed_lines.append(' ' * total_spaces + stripped_line)
                                    else:
                                        processed_lines.append('')
                                else:
                                    processed_lines.append(line.lstrip())

                            content = '\n'.join(processed_lines)
                            logger.info("[txt] 对齐处理完成")

                        logger.debug(f"[txt] LLM 返回内容预览: {content[:300]}...")

                        markers = ['修改后的内容', '修改后内容', '以下是修改后的', '修改后的文档', '文档内容如下',
                                   '修改后的文本']
                        extracted_content = None
                        for marker in markers:
                            if marker in content:
                                idx = content.index(marker)
                                newline_idx = content.find('\n', idx)
                                if newline_idx != -1:
                                    extracted_content = content[newline_idx + 1:].strip()
                                    logger.info(f"[txt] 通过标记 '{marker}' 提取到内容，长度: {len(extracted_content)}")
                                    break

                        if extracted_content is None:
                            if content.startswith('好的') or content.startswith('已经') or content.startswith('我'):
                                paragraphs = content.split('\n\n')
                                if len(paragraphs) > 1:
                                    extracted_content = '\n\n'.join(paragraphs[1:]).strip()
                                    logger.info(f"[txt] 通过段落分割提取到内容，长度: {len(extracted_content)}")

                        final_content = extracted_content if extracted_content else content
                        logger.info(f"[txt] 最终保存内容长度: {len(final_content)}")

                        with open(output_path, 'w', encoding='utf-8') as f:
                            f.write(final_content)

                elif file_ext in ['.xlsx', '.xls']:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path)
                    modified_lines = content.split('\n')
                    current_sheet = None
                    row_idx = 0

                    for line in modified_lines:
                        if line.startswith('=== ') and line.endswith(' ==='):
                            sheet_name = line[4:-4]
                            if sheet_name in wb.sheetnames:
                                current_sheet = wb[sheet_name]
                                row_idx = 1
                            else:
                                current_sheet = wb.create_sheet(sheet_name)
                                row_idx = 1
                        elif current_sheet and line.strip():
                            cells = line.split('\t')
                            for col_idx, cell_value in enumerate(cells, 1):
                                current_sheet.cell(row=row_idx, column=col_idx, value=cell_value)
                            row_idx += 1

                    wb.save(str(output_path))

                elif file_ext in ['.pptx', '.ppt']:
                    from pptx import Presentation
                    prs = Presentation(file_path)

                    pptx_content = content

                    markers = ['修改后的内容', '修改后内容', '以下是修改后的', '修改后的文档']
                    for marker in markers:
                        if marker in content:
                            idx = content.index(marker)
                            newline_idx = content.find('\n', idx)
                            if newline_idx != -1:
                                pptx_content = content[newline_idx + 1:]
                                break

                    modified_lines = pptx_content.split('\n')
                    current_slide_idx = 0
                    current_text_idx = 0

                    for line in modified_lines:
                        if line.startswith('=== 幻灯片 ') and ' ===' in line:
                            slide_num = int(line.split('幻灯片 ')[1].split(' ===')[0])
                            if slide_num <= len(prs.slides):
                                current_slide_idx = slide_num - 1
                                current_text_idx = 0
                        elif line.strip() and current_slide_idx < len(prs.slides):
                            slide = prs.slides[current_slide_idx]
                            text_shapes = [shape for shape in slide.shapes
                                           if hasattr(shape, "has_text_frame") and shape.has_text_frame
                                           and shape.text_frame.text.strip()]
                            if current_text_idx < len(text_shapes):
                                text_shapes[current_text_idx].text_frame.text = line
                                current_text_idx += 1

                    prs.save(str(output_path))

                elif file_ext == '.doc':
                    logger.info("[doc] 将修改后的内容保存为 .docx 格式")
                    docx_output_path = doc_processor.OUTPUT_DIR / f"{original_filename}_{uuid.uuid4().hex[:8]}.docx"

                    from docx import Document
                    doc = Document()

                    modified_lines = content.split('\n')
                    for line in modified_lines:
                        if line.strip():
                            doc.add_paragraph(line)

                    doc.save(str(docx_output_path))

                    try:
                        import win32com.client
                        word = win32com.client.Dispatch("Word.Application")
                        word.Visible = False
                        doc_obj = word.Documents.Open(str(docx_output_path))
                        doc_obj.SaveAs(str(output_path), FileFormat=0)
                        doc_obj.Close()
                        word.Quit()
                        docx_output_path.unlink()
                    except:
                        output_path = docx_output_path
                        output_filename = docx_output_path.name

                if output_path and output_path.exists():
                    download_url = f"http://127.0.0.1:5001/download/{output_filename}"

                    is_typo_request = '错别字' in modification_request or '错字' in modification_request or '错误字' in modification_request

                    if file_ext == '.txt' and is_typo_request:
                        with open(output_path, 'r', encoding='utf-8') as f:
                            modified_content = f.read()

                        original_lines = original_content.strip().split('\n')
                        modified_lines = modified_content.strip().split('\n')

                        changed_words = []
                        for i, (orig_line, mod_line) in enumerate(zip(original_lines, modified_lines)):
                            if orig_line.strip() != mod_line.strip():
                                changed_words.append(orig_line.strip())

                        if changed_words:
                            unique_changes = list(dict.fromkeys(changed_words))
                            changes_str = '，'.join(unique_changes)
                            content = f"发现{len(unique_changes)}组错别字：{changes_str}\n\n文档修改完成！\n\n[下载修改后的文档]({download_url})"
                        else:
                            content = "没有发现错别字"
                    else:
                        content = f"文档修改完成！\n\n[下载修改后的文档]({download_url})"
                else:
                    if content != "没有发现错别字":
                        content = "文档修改完成！"

        except Exception as e:
            logger.exception(f"document_node 处理失败: {str(e)}")
            content = f"抱歉，文档处理失败：{str(e)}"

        return {"messages": [HumanMessage(content=content)], "type": "document", "previous_node": "document_node"}

    def _handle_document_creation(state: State, message_content: str, conversation_history: list) -> dict:
        """处理新建文档逻辑，支持 .docx, .doc, .txt, .xlsx, .xls, .pptx, .ppt 格式"""
        logger.info("document_node - 新建文档模式")

        if '[用户要求]' in message_content:
            parts = message_content.split('[用户要求]')
            creation_request = parts[-1].strip() if len(parts) > 1 else "创建文档"
        else:
            lines = message_content.split('\n')
            creation_request = lines[0].strip() if lines else "创建文档"

        logger.info(f"创建要求: {creation_request}")

        # 检测用户期望的文件格式
        format_keywords = {
            'word': '.docx', 'docx': '.docx', 'doc': '.doc',
            'txt': '.txt', '文本': '.txt', '纯文本': '.txt',
            'excel': '.xlsx', 'xlsx': '.xlsx', 'xls': '.xls', '表格': '.xlsx',
            'ppt': '.pptx', 'pptx': '.pptx', '幻灯片': '.pptx', '演示': '.pptx', 'powerpoint': '.pptx'
        }

        # 首先从当前请求中检测文件类型
        output_ext = None
        for keyword, ext in format_keywords.items():
            if keyword in creation_request.lower():
                output_ext = ext
                break
        
        # 如果当前请求中没有指定文件类型，从上下文历史中推断
        if output_ext is None:
            # 检查最近的对话历史，查找文件类型关键词
            for msg in reversed(conversation_history):
                msg_content = msg.get('content', '') if isinstance(msg, dict) else str(msg)
                for keyword, ext in format_keywords.items():
                    if keyword in msg_content.lower():
                        output_ext = ext
                        logger.info(f'从上下文历史中推断文件类型: {ext} (关键词: {keyword})')
                        break
                if output_ext:
                    break
        
        # 如果仍然没有检测到，使用默认格式
        if output_ext is None:
            output_ext = '.docx'
            logger.info('未检测到文件类型，使用默认格式: .docx')

        system_prompt = f"""你是一个专业的文档创作助手。你的任务是根据用户的要求创建一份完整的文档。

要求：
1. 根据用户需求，生成结构完整、内容丰富的文档
2. 生成结构化文字、分段、标题、列表、表格内容；
3. 调整格式、排版、润色、纠错、逻辑梳理。
4. 内容要专业、准确、有条理
5. 直接返回文档内容，不要任何解释、说明或额外文字
6. 保持清晰的换行结构，便于后续保存为文档文件
7. 输出格式为 {output_ext}

注意：你的回复将直接保存为文档文件，所以只能包含文档内容本身！"""

        try:
            import uuid
            from pathlib import Path

            prompt_list = [
                {"role": "system", "content": system_prompt},
            ] + conversation_history + [
                {"role": "user", "content": f"创建要求：{creation_request}"},
            ]

            response = llm.invoke(prompt_list)
            content = response.content if hasattr(response, 'content') else str(response)

            if not content or content.strip() == "":
                content = "抱歉，暂时无法创建文档。"
            else:
                logger.info(f"LLM 返回内容长度: {len(content)}")

                markers = ['文档内容如下', '以下是文档内容', '创建好的文档', '文档已创建']
                extracted_content = None
                for marker in markers:
                    if marker in content:
                        idx = content.index(marker)
                        newline_idx = content.find('\n', idx)
                        if newline_idx != -1:
                            extracted_content = content[newline_idx + 1:].strip()
                            break

                final_content = extracted_content if extracted_content else content

                doc_type_keywords = {
                    '报告': 'report', '总结': 'summary', '计划': 'plan', '方案': 'proposal',
                    '会议': 'meeting', '记录': 'record', '通知': 'notice', '公告': 'announcement',
                    '合同': 'contract', '协议': 'agreement', '简历': 'resume', '论文': 'paper',
                    '文章': 'article', '演讲稿': 'speech', '策划': 'proposal'
                }

                doc_type = 'document'
                for keyword, type_name in doc_type_keywords.items():
                    if keyword in creation_request:
                        doc_type = type_name
                        break

                output_filename = f"{doc_type}_{uuid.uuid4().hex[:8]}{output_ext}"
                output_path = PathConfig.OUTPUTS_DIR / output_filename

                # 根据文件格式保存
                if output_ext in ['.docx', '.doc']:
                    from docx import Document
                    from docx.shared import Pt
                    import re
                    doc = Document()

                    def clean_md_inline(text):
                        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
                        text = re.sub(r'\*(.+?)\*', r'\1', text)
                        text = re.sub(r'~~(.+?)~~', r'\1', text)
                        text = re.sub(r'`(.+?)`', r'\1', text)
                        return text.strip()

                    def is_table_row(line):
                        s = line.strip()
                        return s.startswith('|') and s.endswith('|') and '|' in s[1:-1]

                    def is_table_separator(line):
                        s = line.strip()
                        # 匹配形如 |---|---|---| 或 |:--|:--|:--| 的分隔行
                        cells = s.strip('|').split('|')
                        return all(re.match(r'^[\s\-:]+$', cell) for cell in cells if cell.strip())

                    def parse_table_rows(lines):
                        if len(lines) < 3:
                            return None, []
                        header = [cell.strip() for cell in lines[0].strip().strip('|').split('|')]
                        data_rows = []
                        for line in lines[2:]:
                            if is_table_row(line) and not is_table_separator(line):
                                row = [cell.strip() for cell in line.strip().strip('|').split('|')]
                                data_rows.append(row)
                        return header, data_rows

                    paragraphs = final_content.split('\n')
                    i = 0
                    while i < len(paragraphs):
                        para = paragraphs[i]

                        # 检测表格
                        if is_table_row(para) and not is_table_separator(para):
                            table_lines = [para]
                            j = i + 1
                            while j < len(paragraphs):
                                if is_table_row(paragraphs[j]) or is_table_separator(paragraphs[j]):
                                    table_lines.append(paragraphs[j])
                                    j += 1
                                else:
                                    break

                            header, data_rows = parse_table_rows(table_lines)
                            if header and data_rows:
                                table = doc.add_table(rows=len(data_rows) + 1, cols=len(header))
                                table.style = 'Table Grid'
                                header_cells = table.rows[0].cells
                                for col_idx, cell_text in enumerate(header):
                                    cell = header_cells[col_idx]
                                    cell.text = clean_md_inline(cell_text)
                                    for paragraph in cell.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.bold = True
                                            run.font.size = Pt(10)
                                for row_idx, row_data in enumerate(data_rows):
                                    row_cells = table.rows[row_idx + 1].cells
                                    for col_idx, cell_text in enumerate(row_data):
                                        if col_idx < len(row_cells):
                                            row_cells[col_idx].text = clean_md_inline(cell_text)
                                            for paragraph in row_cells[col_idx].paragraphs:
                                                for run in paragraph.runs:
                                                    run.font.size = Pt(10)
                                i = j
                                continue
                            else:
                                # 表格解析失败，作为普通段落处理，跳过标题检测
                                doc.add_paragraph(para.strip())
                                i = j
                                continue

                        stripped = para.strip()
                        if not stripped:
                            i += 1
                            continue

                        heading_level = None
                        title_text = None

                        # 1. Markdown 风格标题（绝对级别）
                        if stripped.startswith('#'):
                            title_text = stripped.lstrip('#').strip()
                            if stripped.startswith('######'):
                                heading_level = 6
                            elif stripped.startswith('#####'):
                                heading_level = 5
                            elif stripped.startswith('####'):
                                heading_level = 4
                            elif stripped.startswith('###'):
                                heading_level = 3
                            elif stripped.startswith('##'):
                                heading_level = 2
                            else:
                                heading_level = 1

                        # 2. 中文大写数字标题：一、二、三、四... → 标题2（主章节）
                        elif re.match(r'^[一二三四五六七八九十百]+[、.]', stripped):
                            title_text = stripped
                            heading_level = 2

                        # 3. 阿拉伯数字标题：1. 2. 3. ... → 标题3（子章节）
                        elif re.match(r'^\d+[.、]', stripped):
                            title_text = stripped
                            heading_level = 3

                        # 4. 带括号的数字标题：(1) (2) 或 ① ② → 标题4（孙章节）
                        elif re.match(r'^[\(（]?\d+[\)）]', stripped) or re.match(r'^[①③④⑤⑧⑨⑩]', stripped):
                            title_text = stripped
                            heading_level = 4

                        # 5. 英文字母标题：A. B. C. 或 a. b. c. → 标题4（孙章节）
                        elif re.match(r'^[A-Z][.、]', stripped) or re.match(r'^[a-z][.、]', stripped):
                            title_text = stripped
                            heading_level = 4

                        # 6. 短文本且不以标点结尾 → 标题2（默认章节）
                        elif len(stripped) < 50 and not stripped.endswith(('。', '！', '？', '，', '；', ':', '：')):
                            title_text = stripped
                            heading_level = 2

                        if heading_level and title_text:
                            doc.add_heading(title_text, level=heading_level)
                        else:
                            doc.add_paragraph(stripped)

                        i += 1

                    doc.save(str(output_path))

                elif output_ext == '.txt':
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(final_content)

                elif output_ext in ['.xlsx', '.xls']:
                    import openpyxl
                    from openpyxl.styles import Font, Alignment, Border, Side, PatternFill

                    wb = openpyxl.Workbook()
                    ws = wb.active
                    ws.title = "Sheet1"

                    lines = final_content.split('\n')

                    # 智能解析表格数据
                    parsed_rows = []
                    for line in lines:
                        line = line.strip()
                        if not line:
                            continue

                        # 跳过Markdown表格分隔符行（如 |---|---|）
                        if all(c in '-|: ' for c in line):
                            continue

                        # 解析不同格式的分隔符
                        cells = []
                        if '|' in line:
                            # Markdown表格格式：| 列1 | 列2 | 列3 |
                            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                        elif '\t' in line:
                            # 制表符分隔
                            cells = [cell.strip() for cell in line.split('\t') if cell.strip()]
                        elif ',' in line and line.count(',') > 1:
                            # 逗号分隔（CSV格式）
                            cells = [cell.strip() for cell in line.split(',')]
                        else:
                            # 默认整行作为一列
                            cells = [line]

                        if cells:
                            parsed_rows.append(cells)

                    # 写入Excel
                    for row_idx, row_data in enumerate(parsed_rows, 1):
                        for col_idx, cell_value in enumerate(row_data, 1):
                            cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)

                            # 第一行设置为标题行（加粗、居中、蓝色背景）
                            if row_idx == 1:
                                cell.font = Font(bold=True, color="FFFFFF")
                                cell.alignment = Alignment(horizontal="center", vertical="center")
                                cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                            else:
                                cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)

                    # 自动调整列宽
                    for column in ws.columns:
                        max_length = 0
                        column_letter = column[0].column_letter
                        for cell in column:
                            try:
                                if cell.value:
                                    cell_length = len(str(cell.value))
                                    if cell_length > max_length:
                                        max_length = cell_length
                            except:
                                pass
                        adjusted_width = min(max_length + 2, 50)  # 最大宽度50
                        ws.column_dimensions[column_letter].width = adjusted_width

                    # 添加边框
                    thin_border = Border(
                        left=Side(style='thin'),
                        right=Side(style='thin'),
                        top=Side(style='thin'),
                        bottom=Side(style='thin')
                    )

                    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
                        for cell in row:
                            cell.border = thin_border

                    wb.save(str(output_path))

                elif output_ext in ['.pptx', '.ppt']:
                    from pptx import Presentation
                    from pptx.util import Inches, Pt
                    from pptx.dml.color import RGBColor
                    import re
                    prs = Presentation()
                    prs.slide_width = Inches(13.333)
                    prs.slide_height = Inches(7.5)

                    def clean_md(text):
                        """清理 Markdown 标记"""
                        text = re.sub(r'#{1,6}\s*', '', text)  # 移除 # 标题标记
                        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)  # 移除 ** 加粗
                        text = re.sub(r'\*(.+?)\*', r'\1', text)  # 移除 * 斜体
                        text = re.sub(r'~~(.+?)~~', r'\1', text)  # 移除 ~~ 删除线
                        text = re.sub(r'`(.+?)`', r'\1', text)  # 移除 ` 代码
                        text = re.sub(r'^[-*]\s*', '', text, flags=re.MULTILINE)  # 移除列表标记
                        text = re.sub(r'^\d+[.、]\s*', '', text, flags=re.MULTILINE)  # 移除数字列表
                        return text.strip()

                    def set_font(paragraph, size=18, bold=False, color=None):
                        """设置段落字体"""
                        for run in paragraph.runs:
                            run.font.size = Pt(size)
                            run.font.bold = bold
                            if color:
                                run.font.color.rgb = color

                    # 按幻灯片分隔符分割
                    slide_blocks = re.split(r'\n#{1,3}\s*幻灯片\s*\d+[:：]\s*', final_content)
                    slide_blocks = [b.strip() for b in slide_blocks if b.strip()]

                    for idx, block in enumerate(slide_blocks):
                        lines = block.split('\n')
                        lines = [l.strip() for l in lines if l.strip() and l.strip() != '---']

                        if not lines:
                            continue

                        # 第一行作为标题
                        title_text = clean_md(lines[0])
                        # 移除开头的 "幻灯片X：" 前缀
                        title_text = re.sub(r'^幻灯片\s*\d+[:：]\s*', '', title_text)

                        if idx == 0:
                            # 封面页
                            slide_layout = prs.slide_layouts[6]  # Blank
                            slide = prs.slides.add_slide(slide_layout)

                            # 添加居中标题
                            left = Inches(1)
                            top = Inches(2.5)
                            width = Inches(11.333)
                            height = Inches(1)
                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            tf = txBox.text_frame
                            p = tf.paragraphs[0]
                            p.text = title_text
                            p.alignment = 1  # 居中
                            set_font(p, size=40, bold=True)

                            # 添加副标题（如果有）
                            if len(lines) > 1:
                                top2 = Inches(3.8)
                                txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(2))
                                tf2 = txBox2.text_frame
                                tf2.word_wrap = True
                                for line in lines[1:]:
                                    if line.strip() and line.strip() != '---':
                                        p = tf2.add_paragraph()
                                        p.text = clean_md(line)
                                        p.alignment = 1  # 居中
                                        set_font(p, size=24)
                        else:
                            # 内容页
                            slide_layout = prs.slide_layouts[1]  # Title and Content
                            slide = prs.slides.add_slide(slide_layout)
                            slide.shapes.title.text = title_text

                            # 设置标题字体
                            title_shape = slide.shapes.title
                            for para in title_shape.text_frame.paragraphs:
                                set_font(para, size=32, bold=True)

                            if len(lines) > 1:
                                body_shape = slide.shapes.placeholders[1]
                                tf = body_shape.text_frame
                                tf.clear()

                                first_para = True
                                for line in lines[1:]:
                                    if not line.strip() or line.strip() == '---':
                                        continue

                                    cleaned = clean_md(line)
                                    if not cleaned:
                                        continue

                                    if first_para:
                                        p = tf.paragraphs[0]
                                        first_para = False
                                    else:
                                        p = tf.add_paragraph()

                                    # 检测缩进级别
                                    indent_level = 0
                                    if cleaned.startswith('- ') or cleaned.startswith('• '):
                                        cleaned = cleaned[2:]
                                        indent_level = 0
                                    elif cleaned.startswith('  - ') or cleaned.startswith('    • '):
                                        cleaned = cleaned.lstrip()
                                        if cleaned.startswith('- ') or cleaned.startswith('• '):
                                            cleaned = cleaned[2:]
                                        indent_level = 1

                                    p.text = cleaned
                                    p.level = indent_level
                                    set_font(p, size=20 if indent_level == 0 else 18)

                    prs.save(str(output_path))

                if output_path.exists():
                    download_url = f"http://127.0.0.1:5001/download/{output_filename}"
                    # 返回结构化内容：包含预览内容和下载链接
                    content = f"""文档创建完成！

<document-preview filename="{output_filename}" download-url="{download_url}">
{final_content}
</document-preview>"""
                else:
                    content = "文档创建失败，请重试。"

        except Exception as e:
            logger.exception(f"文档创建失败: {str(e)}")
            content = f"抱歉，文档创建失败：{str(e)}"

        return {"messages": [HumanMessage(content=content)], "type": "document", "previous_node": "document_node"}

    def drawing_node(state: State):
        """图纸识别节点：委托 DrawingAssistant 处理工程图纸相关任务"""
        logger.info("drawing_node 开始处理")

        message = state["messages"][-1]
        if hasattr(message, 'content'):
            user_query = message.content
        else:
            user_query = str(message)

        file_path = state.get("file_path", "")
        assistant = get_drawing_assistant()

        # 根据文件扩展名判断图纸类型
        img_path = None
        pdf_path = None
        dxf_path = None

        if file_path:
            ext = Path(file_path).suffix.lower()
            if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif"):
                img_path = file_path
            elif ext == ".pdf":
                pdf_path = file_path
            elif ext == ".dxf":
                dxf_path = file_path

        try:
            result = assistant.invoke(
                user_query=user_query,
                img_path=img_path,
                pdf_path=pdf_path,
                dxf_path=dxf_path,
                thread_id="agent_draw_session",
            )
        except Exception as e:
            logger.error(f"drawing_node 处理失败: {str(e)}")
            result = f"图纸识别处理失败：{str(e)}"

        return {"messages": [HumanMessage(content=result)], "type": "drawing", "previous_node": "drawing_node"}

    def reflection_node(state: State):
        """反思节点：评估各任务节点的输出质量，决定是否需要重新生成"""
        logger.info("reflection_node 开始反思")

        previous_node = state.get("previous_node", "")
        last_message = state["messages"][-1]
        user_message = state["messages"][-2] if len(state["messages"]) >= 2 else state["messages"][0]

        if hasattr(last_message, 'content'):
            ai_response = last_message.content
        else:
            ai_response = str(last_message)

        if hasattr(user_message, 'content'):
            user_query = user_message.content
        else:
            user_query = str(user_message)

        # 根据上一个节点类型，使用不同的反思标准
        reflection_prompts = {
            "travel_node": """你是一个专业的旅行规划质量评估专家。请评估以下旅行路线规划的质量：

评估标准：
1. 是否包含具体的景点、交通、住宿等实用信息
2. 路线安排是否合理、可行
3. 是否满足用户的具体需求
4. 信息是否准确、详细

用户问题：{user_query}
AI回答：{ai_response}

请给出评估结果，只能返回以下两种之一：
- "PASS"：质量合格，无需修改
- "FAIL: {具体问题描述和改进建议}"：质量不合格，需要重新生成""",

            "joke_node": """你是一个专业的笑话质量评估专家。请评估以下笑话的质量：

评估标准：
1. 是否有趣、有笑点
2. 是否符合用户要求的主题
3. 语言是否流畅自然
4. 长度是否合适（不超过100字）

用户问题：{user_query}
AI回答：{ai_response}

请给出评估结果，只能返回以下两种之一：
- "PASS"：质量合格，无需修改
- "FAIL: {具体问题描述和改进建议}"：质量不合格，需要重新生成""",

            "couplet_node": """你是一个专业的对联质量评估专家。请评估以下对联的质量：

评估标准：
1. 上下联是否对仗工整（词性、平仄、意境）
2. 是否符合用户给出的上联主题
3. 是否有文化内涵和文学美感

用户问题：{user_query}
AI回答：{ai_response}

请给出评估结果，只能返回以下两种之一：
- "PASS"：质量合格，无需修改
- "FAIL: {具体问题描述和改进建议}"：质量不合格，需要重新生成""",

            "document_node": """你是一个专业的文档修改质量评估专家。请评估以下文档修改的质量：

评估标准：
1. 是否按照用户要求进行了修改
2. 修改后的内容是否准确、通顺
3. 是否保留了原文档的核心内容
4. 格式是否正确

用户问题：{user_query}
AI回答：{ai_response}

请给出评估结果，只能返回以下两种之一：
- "PASS"：质量合格，无需修改
- "FAIL: {具体问题描述和改进建议}"：质量不合格，需要重新生成""",

            "code_node": """你是一个专业的代码质量评估专家。请评估以下代码的质量：

评估标准：
1. 代码是否能正确运行（语法正确）
2. 代码逻辑是否合理、无Bug
3. 是否有必要的注释和说明
4. 是否满足用户的具体需求
5. 代码风格是否规范

用户问题：{user_query}
AI回答：{ai_response}

请给出评估结果，只能返回以下两种之一：
- "PASS"：质量合格，无需修改
- "FAIL: {具体问题描述和改进建议}"：质量不合格，需要重新生成""",

            "other_node": """你是一个专业的回答质量评估专家。请评估以下回答的质量：

评估标准：
1. 回答是否礼貌、得体
2. 是否清晰表达了无法处理的原因

用户问题：{user_query}
AI回答：{ai_response}

请给出评估结果，只能返回以下两种之一：
- "PASS"：质量合格，无需修改
- "FAIL: {具体问题描述和改进建议}"：质量不合格，需要重新生成""",

            "drawing_node": """你是一个专业的工程图纸分析质量评估专家。请评估以下图纸分析回答的质量：

评估标准：
1. 是否准确提取了图纸中的尺寸、标注、图层等信息
2. 尺寸分析是否合理、尺寸链是否闭合
3. 国标引用是否准确
4. 加工建议是否专业、实用
5. 回答结构是否清晰（图纸信息、尺寸分析、国标依据、加工建议）

用户问题：{user_query}
AI回答：{ai_response}

请给出评估结果，只能返回以下两种之一：
- "PASS"：质量合格，无需修改
- "FAIL: {具体问题描述和改进建议}"：质量不合格，需要重新生成""",
        }

        reflection_prompt = reflection_prompts.get(previous_node, reflection_prompts["other_node"])

        try:
            response = llm.invoke([
                {"role": "user", "content": reflection_prompt.format(user_query=user_query, ai_response=ai_response)}
            ])
            reflection_result = response.content.strip()
        except Exception as e:
            logger.error(f"reflection_node LLM 调用失败: {str(e)[:50]}")
            reflection_result = "PASS"

        logger.info(f"反思结果: {reflection_result[:50]}...")

        if reflection_result.startswith("PASS"):
            logger.info("反思通过，质量合格")
            return {"type": END}
        else:
            logger.info("反思未通过，需要重新生成")
            feedback = reflection_result.replace("FAIL:", "").strip()
            feedback_message = f"之前的回答质量不够好，请根据以下反馈重新生成：{feedback}"
            return {"messages": [HumanMessage(content=feedback_message)], "type": previous_node.replace("_node", "")}

    def routing_func(state: State):
        current_type = state["type"]
        logger.info(f"routing_func: type = '{current_type}'")
        if current_type == Params.MAPPING_NODE.get(Params.TRAVEL_NODE):
            logger.info(f"路由到 {Params.TRAVEL_NODE}")
            return Params.TRAVEL_NODE
        elif current_type == Params.MAPPING_NODE.get(Params.JOKE_NODE):
            logger.info(f"路由到 {Params.JOKE_NODE}")
            return Params.JOKE_NODE
        elif current_type == Params.MAPPING_NODE.get(Params.COUPLET_NODE):
            logger.info(f"路由到 {Params.COUPLET_NODE}")
            return Params.COUPLET_NODE
        elif current_type == Params.MAPPING_NODE.get(Params.DOCUMENT_NODE):
            logger.info(f"路由到 {Params.DOCUMENT_NODE}")
            return Params.DOCUMENT_NODE
        elif current_type == Params.MAPPING_NODE.get(Params.CODE_NODE):
            logger.info(f"路由到 {Params.CODE_NODE}")
            return Params.CODE_NODE
        elif current_type == Params.MAPPING_NODE.get(Params.DRAWING_NODE):
            logger.info(f"路由到 {Params.DRAWING_NODE}")
            return Params.DRAWING_NODE
        elif current_type == END:
            logger.info("路由到 END")
            return END
        else:
            logger.info(f"路由到 {Params.OTHER_NODE} (未知类型: {current_type})")
            return Params.OTHER_NODE

    def reflection_routing_func(state: State):
        """反思节点的路由：PASS则结束，FAIL则回到原节点重新生成"""
        current_type = state["type"]
        logger.info(f"reflection_routing_func: type = '{current_type}'")
        if current_type == END:
            logger.info("反思通过，路由到 END")
            return END
        else:
            logger.info(f"反思未通过，路由回 {current_type}_node")
            return f"{current_type}_node"

    builder = StateGraph(State)
    builder.add_node(Params.SUPERVISOR_NODE, supervisor_node)
    builder.add_node(Params.TRAVEL_NODE, travel_node)
    builder.add_node(Params.JOKE_NODE, joke_node)
    builder.add_node(Params.COUPLET_NODE, couplet_node)
    builder.add_node(Params.CODE_NODE, code_node)
    builder.add_node(Params.DOCUMENT_NODE, document_node)
    builder.add_node(Params.DRAWING_NODE, drawing_node)
    builder.add_node(Params.OTHER_NODE, other_node)
    builder.add_node(Params.REFLECTION_NODE, reflection_node)
    builder.add_edge(START, Params.SUPERVISOR_NODE)
    builder.add_conditional_edges(Params.SUPERVISOR_NODE, routing_func,
                                  [Params.TRAVEL_NODE, Params.JOKE_NODE, Params.COUPLET_NODE, Params.DOCUMENT_NODE, Params.CODE_NODE, Params.DRAWING_NODE, Params.OTHER_NODE, END])
    builder.add_edge(Params.TRAVEL_NODE, Params.REFLECTION_NODE)
    builder.add_edge(Params.JOKE_NODE, Params.REFLECTION_NODE)
    builder.add_edge(Params.COUPLET_NODE, Params.REFLECTION_NODE)
    builder.add_edge(Params.DOCUMENT_NODE, Params.REFLECTION_NODE)
    builder.add_edge(Params.CODE_NODE, Params.REFLECTION_NODE)
    builder.add_edge(Params.DRAWING_NODE, Params.REFLECTION_NODE)
    builder.add_edge(Params.OTHER_NODE, Params.REFLECTION_NODE)
    builder.add_conditional_edges(Params.REFLECTION_NODE, reflection_routing_func,
                                  [Params.TRAVEL_NODE, Params.JOKE_NODE, Params.COUPLET_NODE, Params.DOCUMENT_NODE, Params.CODE_NODE, Params.DRAWING_NODE, Params.OTHER_NODE, END])
    return builder.compile(checkpointer=InMemorySaver())