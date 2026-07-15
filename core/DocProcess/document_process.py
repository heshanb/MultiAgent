import os
import uuid
import re
from typing import Optional, Dict, Any, Tuple, List
from pathlib import Path
from settings.Define import Params, PathConfig

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

try:
    from unstructured.partition.pdf import partition_pdf
    from unstructured.chunking.title import chunk_by_title
    from unstructured.staging.base import elements_to_markdown
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image
    import pdf2image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

class DocumentProcessor:
    """文档处理工具类，支持多种文档格式的读取、修改和保存"""
    SUPPORTED_FORMATS = Params.DOC_SUPPORTED_FORMATS | {'.md', '.pdf'}
    _BASE_DIR = PathConfig.BASE_DIR
    UPLOAD_DIR = PathConfig.UPLOADS_DIR
    OUTPUT_DIR = PathConfig.OUTPUTS_DIR
    
    def __init__(self):
        self.UPLOAD_DIR.mkdir(exist_ok=True)
        self.OUTPUT_DIR.mkdir(exist_ok=True)
    
    def save_uploaded_file(self, file_content: bytes, filename: str) -> tuple:
        """保存上传的文件并返回 (文件路径, 保存的文件名)"""
        saved_filename = filename
        file_path = self.UPLOAD_DIR / saved_filename
        
        with open(file_path, 'wb') as f:
            f.write(file_content)
        
        return str(file_path), saved_filename
    
    def read_document(self, file_path: str) -> str:
        """读取文档内容并返回文本（统一转换为MD格式）"""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.docx':
            return self._read_docx_to_md(file_path)
        elif ext == '.txt':
            return self._read_txt(file_path)
        elif ext in ['.xlsx', '.xls']:
            return self._read_excel(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self._read_pptx(file_path)
        elif ext in ['.html', '.htm']:
            return self._read_html(file_path)
        elif ext in ['.py', '.java', '.cpp', '.c', '.h', '.js', '.ts', '.jsx', '.tsx', '.css', '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.scala', '.cs', '.sh', '.bash', '.sql', '.yaml', '.yml', '.json', '.xml', '.md', '.ini', '.cfg', '.conf', '.toml']:
            return self._read_txt(file_path)
        elif ext == '.doc':
            return self._read_doc(file_path)
        elif ext == '.pdf':
            return self._read_pdf(file_path)
        else:
            raise ValueError(f"不支持的文件格式: {ext}")
    
    def process_document_modification(self, file_path: str, modification_request: str, 
                                     conversation_history: List[Dict[str, str]], llm) -> Tuple[str, Optional[str]]:
        """处理文档修改请求"""
        ext = Path(file_path).suffix.lower()
        original_content = ""
        
        try:
            if ext == '.docx':
                text_only, text_paragraphs_info = self._read_docx_with_paragraphs(file_path)
            elif ext == '.txt':
                text_only = self._read_txt(file_path)
                original_content = text_only
                text_paragraphs_info = None
            elif ext in ['.xlsx', '.xls']:
                text_only = self._read_excel(file_path)
                text_paragraphs_info = None
            elif ext in ['.pptx', '.ppt']:
                text_only = self._read_pptx(file_path)
                text_paragraphs_info = None
            elif ext == '.doc':
                text_only = self._read_doc(file_path)
                text_paragraphs_info = None
            else:
                raise ValueError(f"不支持的文件格式: {ext}")
            
            system_prompt = self._build_system_prompt(modification_request, ext)
            
            prompt_list = [
                {"role": "system", "content": system_prompt},
            ] + conversation_history + [
                {"role": "user", "content": f"修改要求：{modification_request}\n\n文档内容：\n{text_only}"},
            ]
            
            response = llm.invoke(prompt_list)
            content = response.content if hasattr(response, 'content') else str(response)
            
            if not content or content.strip() == "":
                return "抱歉，暂时无法处理文档修改请求。", None
            
            original_filename = Path(file_path).stem
            output_filename = f"{original_filename}_{uuid.uuid4().hex[:8]}{ext}"
            output_path = self.OUTPUT_DIR / output_filename
            
            if ext == '.docx':
                self._modify_docx_with_format(file_path, content, modification_request, str(output_path), text_paragraphs_info)
            elif ext == '.txt':
                self._modify_txt_with_format(content, str(output_path), modification_request, original_content)
            elif ext in ['.xlsx', '.xls']:
                self._modify_excel(file_path, content, str(output_path))
            elif ext in ['.pptx', '.ppt']:
                self._modify_pptx(file_path, content, str(output_path))
            elif ext == '.doc':
                self._modify_doc(content, str(output_path), original_filename)
            
            if output_path.exists():
                download_url = f"http://127.0.0.1:5001/download/{output_filename}"
                
                if ext == '.txt' and ('错别字' in modification_request or '错字' in modification_request or '错误字' in modification_request):
                    with open(output_path, 'r', encoding='utf-8') as f:
                        modified_content = f.read()
                    
                    original_lines = original_content.strip().split('\n')
                    modified_lines = modified_content.strip().split('\n')
                    
                    changed_words = []
                    for i, (orig_line, mod_line) in enumerate(zip(original_lines, modified_lines)):
                        if orig_line.strip() != mod_line.strip():
                            changed_words.append(orig_line.strip())
                    
                    if changed_words:
                        unique_changes = list(dict.fromkeys(changed_words))
                        changes_str = '，'.join(unique_changes)
                        result_content = f"发现{len(unique_changes)}组错别字：{changes_str}\n\n文档修改完成！\n\n[下载修改后的文档]({download_url})"
                    else:
                        result_content = "没有发现错别字"
                else:
                    result_content = f"文档修改完成！\n\n[下载修改后的文档]({download_url})"
                
                return result_content, str(output_path)
            else:
                if content != "没有发现错别字":
                    return "文档修改完成！", None
                return content, None
                
        except Exception as e:
            return f"抱歉，文档处理失败：{str(e)}", None
    
    def create_document(self, creation_request: str, conversation_history: List[Dict[str, str]], 
                       llm, output_ext: str = '.docx') -> Tuple[str, Optional[str]]:
        """创建新文档"""
        system_prompt = f"""你是一个专业的文档创作助手。你的任务是根据用户的要求创建一份完整的文档。

要求：
1. 根据用户需求，生成结构完整、内容丰富的文档
2. 生成结构化文字、分段、标题、列表、表格内容；
3. 调整格式、排版、润色、纠错、逻辑梳理。
4. 内容要专业、准确、有条理
5. 直接返回文档内容，不要任何解释、说明或额外文字
6. 保持清晰的换行结构，便于后续保存为文档文件
7. 输出格式为 {output_ext}

注意：你的回复将直接保存为文档文件，所以只能包含文档内容本身！"""
        
        try:
            prompt_list = [
                {"role": "system", "content": system_prompt},
            ] + conversation_history + [
                {"role": "user", "content": f"创建要求：{creation_request}"},
            ]
            
            response = llm.invoke(prompt_list)
            content = response.content if hasattr(response, 'content') else str(response)
            
            if not content or content.strip() == "":
                return "抱歉，暂时无法创建文档。", None
            
            markers = ['文档内容如下', '以下是文档内容', '创建好的文档', '文档已创建']
            extracted_content = None
            for marker in markers:
                if marker in content:
                    idx = content.index(marker)
                    newline_idx = content.find('\n', idx)
                    if newline_idx != -1:
                        extracted_content = content[newline_idx + 1:].strip()
                        break
            
            final_content = extracted_content if extracted_content else content
            
            doc_type_keywords = {
                '报告': 'report', '总结': 'summary', '计划': 'plan', '方案': 'proposal',
                '会议': 'meeting', '记录': 'record', '通知': 'notice', '公告': 'announcement',
                '合同': 'contract', '协议': 'agreement', '简历': 'resume', '论文': 'paper',
                '文章': 'article', '演讲稿': 'speech', '策划': 'proposal'
            }
            
            doc_type = 'document'
            for keyword, type_name in doc_type_keywords.items():
                if keyword in creation_request:
                    doc_type = type_name
                    break
            
            output_filename = f"{doc_type}_{uuid.uuid4().hex[:8]}{output_ext}"
            output_path = self.OUTPUT_DIR / output_filename
            
            if output_ext in ['.docx', '.doc']:
                self._create_docx(final_content, str(output_path))
            elif output_ext == '.txt':
                self._create_txt(final_content, str(output_path))
            elif output_ext in ['.xlsx', '.xls']:
                self._create_excel(final_content, str(output_path))
            elif output_ext in ['.pptx', '.ppt']:
                self._create_pptx(final_content, str(output_path))
            
            if output_path.exists():
                download_url = f"http://127.0.0.1:5001/download/{output_filename}"
                result_content = f"""文档创建完成！

<document-preview filename="{output_filename}" download-url="{download_url}">
{final_content}
</document-preview>"""
                return result_content, str(output_path)
            else:
                return "文档创建失败，请重试。", None
                
        except Exception as e:
            return f"抱歉，文档创建失败：{str(e)}", None
    
    def _build_system_prompt(self, modification_request: str, file_ext: str) -> str:
        """根据修改要求构建系统提示词"""
        if '错别字' in modification_request or '错字' in modification_request or '错误字' in modification_request:
            return """你是一个专业的文档校对助手。你的任务是检查并修正文档中的错别字和词语搭配错误。

重要要求：
1. 只返回修改后的文档内容，不要任何解释、说明或额外文字
2. 保持原文档的换行结构和格式完全一致
3. 检查词语搭配是否合理（如"天弃"应改为"天气"，"银绗"应改为"银行"）
4. 修正所有不合理的词语组合和错别字
5. 如果没有发现错别字，直接回复"没有发现错别字"

注意：你的回复将直接保存为文档文件，所以只能包含修改后的文档内容本身！"""
        elif '格式' in modification_request or '排版' in modification_request or '对齐' in modification_request:
            is_alignment_request = '对齐' in modification_request
            
            if is_alignment_request and file_ext == '.txt':
                return """你是一个专业的文档格式调整助手。你的任务是根据用户要求调整文档的对齐方式。

重要要求：
1. 只调整对齐方式（如左对齐、居中对齐、右对齐），不要修改文字内容
2. 对于 .txt 纯文本文件：
   - 左对齐：去除每行开头的所有空格和制表符，让内容紧贴左侧
   - 居中对齐：在每行前后添加适当空格使内容居中
   - 右对齐：在每行开头添加适当空格使内容靠右
3. 保留原文档的所有文字内容
4. 直接返回修改后的文本，保持原文本的换行结构

注意：只调整对齐方式，不要改变任何文字内容！"""
            else:
                return """你是一个专业的文档格式调整助手。你的任务是根据用户要求调整文档格式。

要求：
1. 只调整格式（排版、对齐方式、字体样式等）
2. 不要修改文档的文字内容
3. 保留原文档的所有内容
4. 直接返回修改后的文本，保持原文本的换行结构

注意：只调整格式，不要改变文字内容。"""
        elif '内容' in modification_request or '修改' in modification_request or '增删' in modification_request:
            return """你是一个专业的文档修改助手。你的任务是根据用户要求修改文档内容。

要求：
1. 只修改用户指定的内容部分
2. 保留文档的其他内容不变
3. 保留原文档的结构和格式
4. 直接返回修改后的文本，保持原文本的换行结构

注意：只修改用户要求的部分，不要影响其他内容。"""
        elif '结构' in modification_request or '优化' in modification_request:
            return """你是一个专业的文档优化助手。你的任务是优化文档结构和表达。

要求：
1. 只优化文档结构和表达方式
2. 保留文档的核心内容
3. 不要删除重要信息
4. 直接返回优化后的文本，保持原文本的换行结构

注意：只优化结构和表达，不要改变核心内容。"""
        else:
            return """你是一个专业的文档修改助手。你需要根据用户的要求来修改文档内容。

要求：
1. 只修改用户指定的部分
2. 保留文档的其他内容不变
3. 保留原文档的结构和格式
4. 直接返回修改后的文本，保持原文本的换行结构

注意：严格按照用户的要求进行修改，不要影响其他部分。"""
    
    def _read_docx(self, file_path: str) -> str:
        """读取 .docx 文件"""
        try:
            from docx import Document
            doc = Document(file_path)
            content = []
            for para in doc.paragraphs:
                content.append(para.text)
            return '\n'.join(content)
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")
    
    def _read_docx_to_md(self, file_path: str) -> str:
        """读取 .docx 文件并转换为 Markdown 格式"""
        try:
            from docx import Document
            doc = Document(file_path)
            md_lines = []
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                # 根据段落样式确定标题级别
                style_name = para.style.name.lower()
                if style_name.startswith('heading 1'):
                    md_lines.append(f"# {text}")
                elif style_name.startswith('heading 2'):
                    md_lines.append(f"## {text}")
                elif style_name.startswith('heading 3'):
                    md_lines.append(f"### {text}")
                elif style_name.startswith('heading 4'):
                    md_lines.append(f"#### {text}")
                elif style_name.startswith('heading 5'):
                    md_lines.append(f"##### {text}")
                elif style_name.startswith('heading 6'):
                    md_lines.append(f"###### {text}")
                else:
                    md_lines.append(text)
            
            return '\n\n'.join(md_lines)
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")
    
    def _read_docx_with_paragraphs(self, file_path: str) -> Tuple[str, List[Tuple[int, str]]]:
        """读取 .docx 文件并返回文本和段落信息"""
        from docx import Document
        doc = Document(file_path)
        text_paragraphs_info = []
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                text_paragraphs_info.append((i, para.text))
        text_only = '\n'.join([text for _, text in text_paragraphs_info])
        return text_only, text_paragraphs_info
    
    def _read_txt(self, file_path: str) -> str:
        """读取 .txt 文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _read_html(self, file_path: str) -> str:
        """读取 HTML 文件，提取纯文本内容并转换为 Markdown 格式"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # 使用正则表达式去除HTML标签，提取纯文本
            import re
            
            # 去除脚本和样式标签及其内容
            html_content = re.sub(r'<script[^>]*>[\s\S]*?</script>', '', html_content, flags=re.IGNORECASE)
            html_content = re.sub(r'<style[^>]*>[\s\S]*?</style>', '', html_content, flags=re.IGNORECASE)
            
            # 去除HTML标签，保留文本内容
            text = re.sub(r'<[^>]+>', '', html_content)
            
            # 去除多余的空白字符
            text = re.sub(r'\s+', ' ', text).strip()
            
            # 处理HTML实体
            text = text.replace('&nbsp;', ' ').replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
            text = text.replace('&quot;', '"').replace('&apos;', "'")
            
            return text
        except Exception as e:
            # 如果HTML解析失败，尝试作为普通文本读取
            print(f"HTML解析失败，尝试作为文本文件读取: {str(e)}")
            return self._read_txt(file_path)
    
    def _read_excel(self, file_path: str) -> str:
        """读取 Excel 文件并转换为 Markdown 格式（支持 .xls 和 .xlsx）"""
        ext = Path(file_path).suffix.lower()
        
        # 处理旧版 .xls 格式
        if ext == '.xls':
            return self._read_excel_xls(file_path)
        
        # 处理 .xlsx 格式
        elif ext == '.xlsx':
            return self._read_excel_xlsx(file_path)
        
        else:
            raise ValueError(f"不支持的 Excel 格式: {ext}")
    
    def _read_excel_xls(self, file_path: str) -> str:
        """读取 .xls 格式 Excel 文件"""
        if not XLRD_AVAILABLE:
            raise ImportError("请安装 xlrd: pip install xlrd")
        
        try:
            wb = xlrd.open_workbook(file_path)
            content = []
            for sheet_name in wb.sheet_names():
                ws = wb.sheet_by_name(sheet_name)
                content.append(f"## {sheet_name}")
                
                # 获取所有行数据
                rows = []
                for row_idx in range(ws.nrows):
                    row_data = []
                    for col_idx in range(ws.ncols):
                        cell_value = ws.cell_value(row_idx, col_idx)
                        # 处理 xlrd 的特殊值
                        if isinstance(cell_value, float) and cell_value == int(cell_value):
                            row_data.append(str(int(cell_value)))
                        else:
                            row_data.append(str(cell_value))
                    rows.append(row_data)
                
                if rows:
                    # 转换为MD表格
                    header = '| ' + ' | '.join(rows[0]) + ' |'
                    content.append(header)
                    separator = '| ' + ' | '.join(['---'] * len(rows[0])) + ' |'
                    content.append(separator)
                    for row in rows[1:]:
                        content.append('| ' + ' | '.join(row) + ' |')
                
                content.append('')  # 空行分隔不同sheet
            
            return '\n'.join(content)
        except Exception as e:
            # xlrd读取失败，尝试作为文本文件读取（可能是CSV或其他文本格式）
            print(f"xlrd 读取失败，尝试作为文本文件读取: {str(e)}")
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                # 尝试解析为CSV格式
                return self._parse_text_as_csv(content, file_path)
            except Exception as text_e:
                raise ValueError(f"无法读取文件 '{file_path}'。文件可能不是有效的 Excel 格式，或已损坏。尝试的错误: xlrd: {str(e)}, 文本读取: {str(text_e)}")
    
    def _parse_text_as_csv(self, content: str, file_path: str) -> str:
        """将文本内容解析为CSV格式并转换为MD表格"""
        lines = content.strip().split('\n')
        if not lines:
            return f"## {Path(file_path).name}\n\n文件内容为空或无法解析"
        
        # 尝试识别分隔符
        separators = [',', '\t', ';', '|']
        best_separator = ','
        max_count = 0
        
        for sep in separators:
            count = lines[0].count(sep)
            if count > max_count:
                max_count = count
                best_separator = sep
        
        # 解析CSV
        rows = []
        for line in lines:
            # 跳过空行和BOM
            line = line.strip().replace('\ufeff', '')
            if line:
                rows.append([cell.strip() for cell in line.split(best_separator)])
        
        if rows:
            content = [f"## {Path(file_path).name}"]
            # 转换为MD表格
            header = '| ' + ' | '.join(rows[0]) + ' |'
            content.append(header)
            separator = '| ' + ' | '.join(['---'] * len(rows[0])) + ' |'
            content.append(separator)
            for row in rows[1:]:
                content.append('| ' + ' | '.join(row) + ' |')
            return '\n'.join(content)
        
        return f"## {Path(file_path).name}\n\n{content}"
    
    def _read_excel_xlsx(self, file_path: str) -> str:
        """读取 .xlsx 格式 Excel 文件"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            content = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                content.append(f"## {sheet}")
                
                # 获取所有行数据
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(cell) if cell is not None else '' for cell in row])
                
                if rows:
                    # 转换为MD表格
                    header = '| ' + ' | '.join(rows[0]) + ' |'
                    content.append(header)
                    separator = '| ' + ' | '.join(['---'] * len(rows[0])) + ' |'
                    content.append(separator)
                    for row in rows[1:]:
                        content.append('| ' + ' | '.join(row) + ' |')
                
                content.append('')  # 空行分隔不同sheet
            
            return '\n'.join(content)
        except ImportError:
            raise ImportError("请安装 openpyxl: pip install openpyxl")
    
    def _read_pptx(self, file_path: str) -> str:
        """读取 PowerPoint 文件并转换为 Markdown 格式"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content_parts = []

            for i, slide in enumerate(prs.slides, 1):
                content_parts.append(f"## 幻灯片 {i}")
                
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        if shape.text_frame and hasattr(shape.text_frame, 'text'):
                            text_content = shape.text_frame.text.strip()
                            if text_content:
                                slide_texts.append(text_content)
                
                if slide_texts:
                    # 尝试识别标题（通常是第一张形状或最大的字体）
                    if slide_texts:
                        content_parts.append(f"### {slide_texts[0]}")
                        for text in slide_texts[1:]:
                            content_parts.append(text)
                
                content_parts.append('')  # 空行分隔不同幻灯片

            return '\n'.join(content_parts)
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")
    
    def _read_doc(self, file_path: str) -> str:
        """读取旧版 .doc 文件"""
        try:
            import win32com.client
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc_obj = word.Documents.Open(file_path)
            text_only = doc_obj.Content.Text
            doc_obj.Close()
            word.Quit()
            return text_only
        except ImportError:
            raise ValueError("暂不支持直接读取 .doc 文件，请先转换为 .docx 格式")
    
    def _read_pdf(self, file_path: str) -> str:
        """读取 PDF 文件内容（支持扫描件OCR识别，并转换为MD格式）"""
        # 优先使用 unstructured 库进行智能分区解析并转换为MD格式
        if UNSTRUCTURED_AVAILABLE:
            try:
                print(f"使用 unstructured 解析 PDF: {file_path}")
                elements = partition_pdf(file_path)
                
                if not elements:
                    print("unstructured 未解析到任何内容")
                else:
                    print(f"unstructured 解析到 {len(elements)} 个元素")
                
                # 将元素转换为 Markdown 格式（保留标题、表格等结构）
                text = elements_to_markdown(elements)
                print(f"unstructured 提取总文本长度: {len(text)}")
                
                if text.strip():
                    return text.strip()
            except Exception as e:
                print(f"unstructured 解析失败，尝试使用 pypdf: {str(e)}")
        
        # 回退到 pypdf
        if PdfReader is not None:
            try:
                with open(file_path, 'rb') as f:
                    reader = PdfReader(f)
                    text = ""
                    num_pages = len(reader.pages)
                    print(f"PDF 文件页数: {num_pages}")
                    
                    for i, page in enumerate(reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            print(f"第 {i+1} 页提取到文本长度: {len(page_text)}")
                            text += page_text + "\n\n"
                        else:
                            print(f"第 {i+1} 页未提取到文本")
                    
                    print(f"pypdf 提取总文本长度: {len(text)}")
                    if text.strip():
                        return text.strip()
            except Exception as e:
                print(f"pypdf 读取错误: {str(e)}")
        
        # 如果以上方法都失败，尝试OCR（处理扫描件）
        if OCR_AVAILABLE:
            try:
                print(f"尝试使用 OCR 解析扫描件 PDF: {file_path}")
                # 将PDF转换为图像
                images = pdf2image.convert_from_path(file_path)
                print(f"PDF 转换为 {len(images)} 张图片")
                
                text = ""
                for i, image in enumerate(images):
                    # 使用Tesseract进行OCR识别
                    page_text = pytesseract.image_to_string(image, lang='chi_sim')
                    print(f"OCR 第 {i+1} 页提取到文本长度: {len(page_text)}")
                    if page_text:
                        text += page_text + "\n\n"
                
                print(f"OCR 提取总文本长度: {len(text)}")
                if text.strip():
                    return text.strip()
            except Exception as e:
                print(f"OCR 解析失败: {str(e)}")
        
        # 如果所有方法都失败
        raise ValueError(f"无法从 PDF 文件中提取文本，可能是扫描件或加密文件")
    
    def _modify_docx_with_format(self, original_path: str, modified_content: str, 
                                instructions: str, output_path: str, paragraph_info: List[Tuple[int, str]]):
        """修改 .docx 文件并保持格式"""
        from docx import Document
        doc = Document(original_path)
        modified_lines = modified_content.split('\n')
        
        for idx, (para_index, original_text) in enumerate(paragraph_info):
            if idx < len(modified_lines):
                para = doc.paragraphs[para_index]
                new_text = modified_lines[idx]
                for run in para.runs:
                    run.text = ""
                if new_text:
                    para.add_run(new_text)
        
        doc.save(output_path)
    
    def _modify_txt_with_format(self, modified_content: str, output_path: str, 
                               instructions: str, original_content: str):
        """修改 .txt 文件并处理格式"""
        is_typo_request = '错别字' in instructions or '错字' in instructions or '错误字' in instructions
        
        if is_typo_request:
            no_change_markers = ['没有发现', '未发现', '没有错别字', '无需修改', '原文档内容', '文档内容正确']
            is_no_change = any(marker in modified_content for marker in no_change_markers)
            
            if is_no_change:
                return
            
            markers = ['修改后的内容', '修改后内容', '以下是修改后的', '修改后的文档', '文档内容如下', '修改后的文本']
            extracted_content = None
            for marker in markers:
                if marker in modified_content:
                    idx = modified_content.index(marker)
                    newline_idx = modified_content.find('\n', idx)
                    if newline_idx != -1:
                        extracted_content = modified_content[newline_idx + 1:].strip()
                        break
            
            if extracted_content is None:
                if modified_content.startswith('好的') or modified_content.startswith('已经') or modified_content.startswith('我'):
                    paragraphs = modified_content.split('\n\n')
                    if len(paragraphs) > 1:
                        extracted_content = '\n\n'.join(paragraphs[1:]).strip()
            
            final_content = extracted_content if extracted_content else modified_content
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(final_content)
        else:
            if '对齐' in instructions or '格式' in instructions:
                lines = modified_content.split('\n')
                processed_lines = []
                
                for line in lines:
                    if '左对齐' in instructions or '左' in instructions:
                        processed_lines.append(line.lstrip())
                    elif '居中' in instructions:
                        stripped_line = line.strip()
                        if stripped_line:
                            total_spaces = max(0, (80 - len(stripped_line)) // 2)
                            processed_lines.append(' ' * total_spaces + stripped_line)
                        else:
                            processed_lines.append('')
                    elif '右对齐' in instructions:
                        stripped_line = line.strip()
                        if stripped_line:
                            total_spaces = max(0, 80 - len(stripped_line))
                            processed_lines.append(' ' * total_spaces + stripped_line)
                        else:
                            processed_lines.append('')
                    else:
                        processed_lines.append(line.lstrip())
                
                modified_content = '\n'.join(processed_lines)
            
            markers = ['修改后的内容', '修改后内容', '以下是修改后的', '修改后的文档', '文档内容如下', '修改后的文本']
            extracted_content = None
            for marker in markers:
                if marker in modified_content:
                    idx = modified_content.index(marker)
                    newline_idx = modified_content.find('\n', idx)
                    if newline_idx != -1:
                        extracted_content = modified_content[newline_idx + 1:].strip()
                        break
            
            if extracted_content is None:
                if modified_content.startswith('好的') or modified_content.startswith('已经') or modified_content.startswith('我'):
                    paragraphs = modified_content.split('\n\n')
                    if len(paragraphs) > 1:
                        extracted_content = '\n\n'.join(paragraphs[1:]).strip()
            
            final_content = extracted_content if extracted_content else modified_content
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(final_content)
    
    def _modify_excel(self, original_path: str, modified_content: str, output_path: str):
        """修改 Excel 文件"""
        import openpyxl
        wb = openpyxl.load_workbook(original_path)
        modified_lines = modified_content.split('\n')
        current_sheet = None
        row_idx = 0
        
        for line in modified_lines:
            if line.startswith('=== ') and line.endswith(' ==='):
                sheet_name = line[4:-4]
                if sheet_name in wb.sheetnames:
                    current_sheet = wb[sheet_name]
                    row_idx = 1
                else:
                    current_sheet = wb.create_sheet(sheet_name)
                    row_idx = 1
            elif current_sheet and line.strip():
                cells = line.split('\t')
                for col_idx, cell_value in enumerate(cells, 1):
                    current_sheet.cell(row=row_idx, column=col_idx, value=cell_value)
                row_idx += 1
        
        wb.save(output_path)
    
    def _modify_pptx(self, original_path: str, modified_content: str, output_path: str):
        """修改 PowerPoint 文件"""
        from pptx import Presentation
        prs = Presentation(original_path)
        
        pptx_content = modified_content
        
        markers = ['修改后的内容', '修改后内容', '以下是修改后的', '修改后的文档']
        for marker in markers:
            if marker in modified_content:
                idx = modified_content.index(marker)
                newline_idx = modified_content.find('\n', idx)
                if newline_idx != -1:
                    pptx_content = modified_content[newline_idx + 1:]
                    break
        
        modified_lines = pptx_content.split('\n')
        current_slide_idx = 0
        current_text_idx = 0
        
        for line in modified_lines:
            if line.startswith('=== 幻灯片 ') and ' ===' in line:
                slide_num = int(line.split('幻灯片 ')[1].split(' ===')[0])
                if slide_num <= len(prs.slides):
                    current_slide_idx = slide_num - 1
                    current_text_idx = 0
            elif line.strip() and current_slide_idx < len(prs.slides):
                slide = prs.slides[current_slide_idx]
                text_shapes = [shape for shape in slide.shapes
                               if hasattr(shape, "has_text_frame") and shape.has_text_frame
                               and shape.text_frame.text.strip()]
                if current_text_idx < len(text_shapes):
                    text_shapes[current_text_idx].text_frame.text = line
                    current_text_idx += 1
        
        prs.save(output_path)
    
    def _modify_doc(self, modified_content: str, output_path: str, original_filename: str):
        """修改 .doc 文件（转换为 .docx）"""
        docx_output_path = self.OUTPUT_DIR / f"{original_filename}_{uuid.uuid4().hex[:8]}.docx"
        
        from docx import Document
        doc = Document()
        
        modified_lines = modified_content.split('\n')
        for line in modified_lines:
            if line.strip():
                doc.add_paragraph(line)
        
        doc.save(str(docx_output_path))
        
        try:
            import win32com.client
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc_obj = word.Documents.Open(str(docx_output_path))
            doc_obj.SaveAs(str(output_path), FileFormat=0)
            doc_obj.Close()
            word.Quit()
            docx_output_path.unlink()
        except:
            os.replace(str(docx_output_path), str(output_path))
    
    def _create_docx(self, content: str, output_path: str):
        """创建 .docx 文档"""
        from docx import Document
        from docx.shared import Pt
        
        doc = Document()
        
        def clean_md_inline(text):
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            text = re.sub(r'\*(.+?)\*', r'\1', text)
            text = re.sub(r'~~(.+?)~~', r'\1', text)
            text = re.sub(r'`(.+?)`', r'\1', text)
            return text.strip()
        
        def is_table_row(line):
            s = line.strip()
            return s.startswith('|') and s.endswith('|') and '|' in s[1:-1]
        
        def is_table_separator(line):
            s = line.strip()
            cells = s.strip('|').split('|')
            return all(re.match(r'^[\s\-:]+$', cell) for cell in cells if cell.strip())
        
        def parse_table_rows(lines):
            if len(lines) < 3:
                return None, []
            header = [cell.strip() for cell in lines[0].strip().strip('|').split('|')]
            data_rows = []
            for line in lines[2:]:
                if is_table_row(line) and not is_table_separator(line):
                    row = [cell.strip() for cell in line.strip().strip('|').split('|')]
                    data_rows.append(row)
            return header, data_rows
        
        paragraphs = content.split('\n')
        i = 0
        while i < len(paragraphs):
            para = paragraphs[i]
            
            if is_table_row(para) and not is_table_separator(para):
                table_lines = [para]
                j = i + 1
                while j < len(paragraphs):
                    if is_table_row(paragraphs[j]) or is_table_separator(paragraphs[j]):
                        table_lines.append(paragraphs[j])
                        j += 1
                    else:
                        break
                
                header, data_rows = parse_table_rows(table_lines)
                if header and data_rows:
                    table = doc.add_table(rows=len(data_rows) + 1, cols=len(header))
                    table.style = 'Table Grid'
                    header_cells = table.rows[0].cells
                    for col_idx, cell_text in enumerate(header):
                        cell = header_cells[col_idx]
                        cell.text = clean_md_inline(cell_text)
                        for paragraph in cell.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True
                                run.font.size = Pt(10)
                    for row_idx, row_data in enumerate(data_rows):
                        row_cells = table.rows[row_idx + 1].cells
                        for col_idx, cell_text in enumerate(row_data):
                            if col_idx < len(row_cells):
                                row_cells[col_idx].text = clean_md_inline(cell_text)
                                for paragraph in row_cells[col_idx].paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(10)
                    i = j
                    continue
                else:
                    doc.add_paragraph(para.strip())
                    i = j
                    continue
            
            stripped = para.strip()
            if not stripped:
                i += 1
                continue
            
            heading_level = None
            title_text = None
            
            if stripped.startswith('#'):
                title_text = stripped.lstrip('#').strip()
                if stripped.startswith('######'):
                    heading_level = 6
                elif stripped.startswith('#####'):
                    heading_level = 5
                elif stripped.startswith('####'):
                    heading_level = 4
                elif stripped.startswith('###'):
                    heading_level = 3
                elif stripped.startswith('##'):
                    heading_level = 2
                else:
                    heading_level = 1
            elif re.match(r'^[一二三四五六七八九十百]+[、.]', stripped):
                title_text = stripped
                heading_level = 2
            elif re.match(r'^\d+[.、]', stripped):
                title_text = stripped
                heading_level = 3
            elif re.match(r'^[\(（]?\d+[\)）]', stripped) or re.match(r'^[①③④⑤⑧⑨⑩]', stripped):
                title_text = stripped
                heading_level = 4
            elif re.match(r'^[A-Z][.、]', stripped) or re.match(r'^[a-z][.、]', stripped):
                title_text = stripped
                heading_level = 4
            elif len(stripped) < 50 and not stripped.endswith(('。', '！', '？', '，', '；', ':', '：')):
                title_text = stripped
                heading_level = 2
            
            if heading_level and title_text:
                doc.add_heading(title_text, level=heading_level)
            else:
                doc.add_paragraph(stripped)
            
            i += 1
        
        doc.save(output_path)
    
    def _create_txt(self, content: str, output_path: str):
        """创建 .txt 文档"""
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
    
    def _create_excel(self, content: str, output_path: str):
        """创建 Excel 文档"""
        import openpyxl
        from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
        
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        
        lines = content.split('\n')
        parsed_rows = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            if all(c in '-|: ' for c in line):
                continue
            
            cells = []
            if '|' in line:
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            elif '\t' in line:
                cells = [cell.strip() for cell in line.split('\t') if cell.strip()]
            elif ',' in line and line.count(',') > 1:
                cells = [cell.strip() for cell in line.split(',')]
            else:
                cells = [line]
            
            if cells:
                parsed_rows.append(cells)
        
        for row_idx, row_data in enumerate(parsed_rows, 1):
            for col_idx, cell_value in enumerate(row_data, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                
                if row_idx == 1:
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.alignment = Alignment(horizontal="center", vertical="center")
                    cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                else:
                    cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if cell.value:
                        cell_length = len(str(cell.value))
                        if cell_length > max_length:
                            max_length = cell_length
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column_letter].width = adjusted_width
        
        thin_border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
        
        for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
            for cell in row:
                cell.border = thin_border
        
        wb.save(output_path)
    
    def _create_pptx(self, content: str, output_path: str):
        """创建 PowerPoint 文档"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        def clean_md(text):
            text = re.sub(r'#{1,6}\s*', '', text)
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            text = re.sub(r'\*(.+?)\*', r'\1', text)
            text = re.sub(r'~~(.+?)~~', r'\1', text)
            text = re.sub(r'`(.+?)`', r'\1', text)
            text = re.sub(r'^[-*]\s*', '', text, flags=re.MULTILINE)
            text = re.sub(r'^\d+[.、]\s*', '', text, flags=re.MULTILINE)
            return text.strip()
        
        def set_font(paragraph, size=18, bold=False, color=None):
            for run in paragraph.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                if color:
                    run.font.color.rgb = color
        
        slide_blocks = re.split(r'\n#{1,3}\s*幻灯片\s*\d+[:：]\s*', content)
        slide_blocks = [b.strip() for b in slide_blocks if b.strip()]
        
        for idx, block in enumerate(slide_blocks):
            lines = block.split('\n')
            lines = [l.strip() for l in lines if l.strip() and l.strip() != '---']
            
            if not lines:
                continue
            
            title_text = clean_md(lines[0])
            title_text = re.sub(r'^幻灯片\s*\d+[:：]\s*', '', title_text)
            
            if idx == 0:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                
                left = Inches(1)
                top = Inches(2.5)
                width = Inches(11.333)
                height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = title_text
                p.alignment = 1
                set_font(p, size=40, bold=True)
                
                if len(lines) > 1:
                    top2 = Inches(3.8)
                    txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(2))
                    tf2 = txBox2.text_frame
                    tf2.word_wrap = True
                    for line in lines[1:]:
                        if line.strip() and line.strip() != '---':
                            p = tf2.add_paragraph()
                            p.text = clean_md(line)
                            p.alignment = 1
                            set_font(p, size=24)
            else:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title_text
                
                title_shape = slide.shapes.title
                for para in title_shape.text_frame.paragraphs:
                    set_font(para, size=32, bold=True)
                
                if len(lines) > 1:
                    body_shape = slide.shapes.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear()
                    
                    first_para = True
                    for line in lines[1:]:
                        if not line.strip() or line.strip() == '---':
                            continue
                        
                        cleaned = clean_md(line)
                        if not cleaned:
                            continue
                        
                        if first_para:
                            p = tf.paragraphs[0]
                            first_para = False
                        else:
                            p = tf.add_paragraph()
                        
                        indent_level = 0
                        if cleaned.startswith('- ') or cleaned.startswith('• '):
                            cleaned = cleaned[2:]
                            indent_level = 0
                        elif cleaned.startswith('  - ') or cleaned.startswith('    • '):
                            cleaned = cleaned.lstrip()
                            if cleaned.startswith('- ') or cleaned.startswith('• '):
                                cleaned = cleaned[2:]
                            indent_level = 1
                        
                        p.text = cleaned
                        p.level = indent_level
                        set_font(p, size=20 if indent_level == 0 else 18)
        
        prs.save(output_path)
    
    def detect_output_format(self, creation_request: str, conversation_history: List[Dict[str, str]]) -> str:
        """检测用户期望的输出文件格式"""
        format_keywords = {
            'word': '.docx', 'docx': '.docx', 'doc': '.doc',
            'txt': '.txt', '文本': '.txt', '纯文本': '.txt',
            'excel': '.xlsx', 'xlsx': '.xlsx', 'xls': '.xls', '表格': '.xlsx',
            'ppt': '.pptx', 'pptx': '.pptx', '幻灯片': '.pptx', '演示': '.pptx', 'powerpoint': '.pptx'
        }
        
        for keyword, ext in format_keywords.items():
            if keyword in creation_request.lower():
                return ext
        
        for msg in reversed(conversation_history):
            msg_content = msg.get('content', '') if isinstance(msg, dict) else str(msg)
            for keyword, ext in format_keywords.items():
                if keyword in msg_content.lower():
                    return ext
        
        return '.docx'
    
    def get_download_url(self, file_path: str) -> str:
        """生成下载链接"""
        return f"/download/{Path(file_path).name}"


doc_processor = DocumentProcessor()